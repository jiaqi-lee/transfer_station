from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


TITLE_COLOR = RGBColor(31, 78, 121)
TEXT_COLOR = RGBColor(35, 35, 35)
ACCENT = RGBColor(91, 155, 213)
HEADER = RGBColor(221, 235, 247)
FONT = "Microsoft YaHei"
OUT = "flashcomm_pass_poc_solution.pptx"


def setup():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    return prs


def set_font(paragraph, size, bold=False, color=TEXT_COLOR):
    paragraph.font.name = FONT
    paragraph.font.size = Pt(size)
    paragraph.font.bold = bold
    paragraph.font.color.rgb = color


def add_title(slide, title, subtitle=None):
    box = slide.shapes.add_textbox(Inches(0.45), Inches(0.25), Inches(12.4), Inches(0.6))
    p = box.text_frame.paragraphs[0]
    p.text = title
    set_font(p, 26, True, TITLE_COLOR)
    if subtitle:
        sub = slide.shapes.add_textbox(Inches(0.48), Inches(0.86), Inches(12.2), Inches(0.35))
        sp = sub.text_frame.paragraphs[0]
        sp.text = subtitle
        set_font(sp, 12, False, RGBColor(100, 100, 100))


def add_bullets(slide, left, top, width, height, bullets, font_size=15):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.08)
    tf.margin_right = Inches(0.08)
    for i, item in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        text, level = item if isinstance(item, tuple) else (item, 0)
        p.text = text
        p.level = level
        set_font(p, font_size if level == 0 else font_size - 1)
        p.space_after = Pt(5)


def add_table(slide, left, top, width, height, data, font_size=11):
    table = slide.shapes.add_table(len(data), len(data[0]), Inches(left), Inches(top), Inches(width), Inches(height)).table
    for r, row in enumerate(data):
        for c, text in enumerate(row):
            cell = table.cell(r, c)
            cell.text = text
            if r == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = HEADER
            for p in cell.text_frame.paragraphs:
                set_font(p, font_size, r == 0)


def add_box(slide, left, top, width, height, text, fill=RGBColor(242, 248, 252), font_size=14):
    shape = slide.shapes.add_shape(1, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = ACCENT
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = PP_ALIGN.CENTER
    set_font(p, font_size, True)
    return shape


def add_arrow(slide, left, top, width, height):
    shape = slide.shapes.add_shape(33, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT
    shape.line.color.rgb = ACCENT
    return shape


def new_slide(prs, title, subtitle=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, title, subtitle)
    return slide


def main():
    prs = setup()

    s = new_slide(prs, "FlashComm Dense Pass POC 方案说明", "用于和 SE 对齐 FX pass 化整改方向")
    add_bullets(s, 0.75, 1.35, 11.9, 5.3, [
        "目标：用一个独立、可开关的 POC 验证 FlashComm v1 基础场景能否通过 FX graph pass 表达。",
        "本次只做 dense 基础 pattern：all_reduce + AddRMSNormBias -> reduce_scatter + AddRMSNormBias + all_gather。",
        "不在 POC 中处理 shared_expert_dp、multistream_moe、多模态第一层、MoE allgather 后移、量化 allgather 后置。",
        "设计原则：先验证 pass 思路是否正确，再拿结果和 SE 讨论复杂特性叠加是否继续扩展更多 pattern。",
    ], 18)

    s = new_slide(prs, "为什么先做这个 POC")
    add_bullets(s, 0.75, 1.25, 12.0, 5.4, [
        "当前 FlashComm v1 逻辑散落在 layer、自定义算子、forward context 和 MoE prepare/finalize 中，if/else 分支多。",
        "需求目标希望对齐主社区 sequence_parallel，用 FX pass 做算子 pattern 替换，减少显式修改 layer 代码。",
        "dense 场景的通信/norm 关系最清晰，适合作为第一阶段穿刺：只证明“图上能匹配并替换”。",
        "如果 dense 基础方案被认可，再讨论 MoE、多流、shared expert、量化等复杂场景如何扩展。",
    ], 16)

    s = new_slide(prs, "POC 覆盖范围与非目标")
    add_table(s, 0.55, 1.25, 12.2, 5.0, [
        ["类别", "本 POC 处理情况", "原因"],
        ["Dense 基础 RMSNorm", "覆盖", "pattern 稳定，等价于 FlashComm/SP 最核心替换"],
        ["Middle layer", "覆盖", "处理 result 和 residual 两个输出"],
        ["Last layer", "覆盖", "只处理 result 输出，不再传播 residual"],
        ["MoE / shared_expert_dp", "不覆盖", "涉及 shared/routed 输出聚合和 prepare/finalize 语义"],
        ["multistream_moe", "不覆盖", "涉及 stream/event/context，非纯算子序列替换"],
        ["多模态第一层", "不覆盖", "缺少 embedding 后 allreduce，pattern 和 dense LLM 不一致"],
        ["量化", "不覆盖", "allgather 与 quant 的相对位置需要单独设计"],
    ], 11)

    s = new_slide(prs, "核心替换 pattern")
    add_box(s, 0.7, 1.45, 3.0, 0.8, "原始图\nAllReduce")
    add_arrow(s, 3.9, 1.63, 0.7, 0.35)
    add_box(s, 4.75, 1.45, 3.4, 0.8, "AddRMSNormBias\n全 token 计算")
    add_arrow(s, 8.35, 1.63, 0.7, 0.35)
    add_box(s, 9.25, 1.45, 3.0, 0.8, "输出\nfull token")
    add_box(s, 0.7, 3.45, 3.0, 0.8, "替换后\nReduceScatter")
    add_arrow(s, 3.9, 3.63, 0.7, 0.35)
    add_box(s, 4.75, 3.45, 3.4, 0.8, "AddRMSNormBias\ntoken shard 计算")
    add_arrow(s, 8.35, 3.63, 0.7, 0.35)
    add_box(s, 9.25, 3.45, 3.0, 0.8, "AllGather\n恢复 full token")
    add_bullets(s, 0.75, 5.25, 11.9, 1.0, [
        "收益方向：norm 只在 token shard 上计算，通信从单个 allreduce 拆成 RS/AG，给后续通信计算融合留下空间。"
    ], 15)

    s = new_slide(prs, "实现结构：新增 FlashCommDensePass")
    add_table(s, 0.55, 1.25, 12.2, 4.6, [
        ["代码位置", "职责"],
        ["vllm_ascend/compilation/passes/flashcomm_pass.py", "定义 FlashCommDensePass"],
        ["PatternMatcherPass(pass_name='npu_flashcomm_dense_pass')", "承载 dense flashcomm pattern"],
        ["MiddleAllReduceRMSNormPattern", "复用现有 middle-layer allreduce+rmsnorm 替换规则"],
        ["LastAllReduceRMSNormPattern", "复用现有 last-layer allreduce+rmsnorm 替换规则"],
        ["NoOpEliminationPass", "替换前先清理 view-like no-op，提升 pattern 命中稳定性"],
        ["get_sp_min_token_num(config)", "复用现有 dense/MoE token 阈值策略"],
    ], 11)
    add_bullets(s, 0.75, 6.15, 11.8, 0.6, ["这版 POC 复用已有成熟 pattern，只把“FlashComm 命名和独立开关”作为验证载体。"], 14)

    s = new_slide(prs, "接入方式：独立开关，避免影响现有 SP")
    add_bullets(s, 0.75, 1.25, 12.0, 2.1, [
        "接入点：GraphFusionPassManager.configure()。",
        "新增开关：additional_config.ascend_compilation_config.enable_flashcomm_pass_poc。",
        "保护条件：只有 pass_config.enable_sp 为 False 时才注册 FlashCommDensePass，避免和现有 SequenceParallelismPass 重复注册同类 pattern。",
    ], 16)
    add_table(s, 0.75, 3.85, 11.8, 2.1, [
        ["配置组合", "注册结果"],
        ["enable_flashcomm_pass_poc=false", "不注册 POC pass"],
        ["enable_flashcomm_pass_poc=true 且 enable_sp=false", "注册 FlashCommDensePass"],
        ["enable_flashcomm_pass_poc=true 且 enable_sp=true", "跳过 POC pass，继续走现有 SequenceParallelismPass + SequenceParallelismMoePass"],
    ], 12)

    s = new_slide(prs, "示例配置")
    add_bullets(s, 0.75, 1.3, 12.0, 1.2, [
        "用于 POC 穿刺时，可通过 additional_config 打开，不改变社区 pass_config.enable_sp 现有路径：",
    ], 16)
    add_box(s, 1.0, 2.45, 11.3, 2.35, '{\n  "ascend_compilation_config": {\n    "enable_flashcomm_pass_poc": true\n  }\n}', RGBColor(245, 245, 245), 16)
    add_bullets(s, 0.75, 5.35, 12.0, 1.0, [
        "注意：这是 POC 开关，不建议直接作为最终产品接口；最终命名应和 SE 确认是否对齐 sequence_parallel。"
    ], 15)

    s = new_slide(prs, "执行流程")
    add_box(s, 0.7, 1.35, 2.2, 0.8, "Dynamo/AOT\n生成 FX 图")
    add_arrow(s, 3.05, 1.57, 0.55, 0.3)
    add_box(s, 3.75, 1.35, 2.4, 0.8, "GraphFusion\nPassManager")
    add_arrow(s, 6.3, 1.57, 0.55, 0.3)
    add_box(s, 7.0, 1.35, 2.5, 0.8, "FlashCommDensePass\n可选注册")
    add_arrow(s, 9.65, 1.57, 0.55, 0.3)
    add_box(s, 10.35, 1.35, 2.2, 0.8, "graph.recompile")
    add_bullets(s, 0.75, 3.05, 12.0, 3.5, [
        "Pass 内部顺序：begin -> NoOpEliminationPass -> PatternMatcherPass.apply -> 记录 matched_count -> end_and_log。",
        "Pattern 使用 torch._inductor.pattern_matcher.register_replacement 注册。",
        "替换后的图仍使用 torch.ops.vllm.reduce_scatter / all_gather 和 torch.ops._C_ascend.npu_add_rms_norm_bias。",
    ], 16)

    s = new_slide(prs, "Middle layer 与 Last layer 的差异")
    add_table(s, 0.55, 1.25, 12.2, 4.9, [
        ["Pattern", "原始返回", "替换后处理", "说明"],
        ["MiddleAllReduceRMSNormPattern", "(result, residual)", "RS 后 chunk residual，再 norm，再 AG result", "中间层 residual 继续参与后续层"],
        ["LastAllReduceRMSNormPattern", "result", "RS 后 chunk residual，再 norm，再 AG result", "最后层只需要输出 result"],
        ["Qwen3VLMiddleAllReduceRMSNormPattern", "未纳入 POC", "不处理 deepstack_input_embeds", "多模态中间层 pattern 单独讨论"],
    ], 11)
    add_bullets(s, 0.75, 6.15, 11.8, 0.6, ["POC 有意不把多模态 pattern 放进来，避免 SE 评审时把基础方案和 VL 特判混在一起。"], 14)

    s = new_slide(prs, "为什么复用现有 pattern，而不是重写一套")
    add_bullets(s, 0.75, 1.25, 12.0, 5.4, [
        "现有 SequenceParallelismPass 已经表达了 Ascend 上的核心 FlashComm/SP 替换语义。",
        "POC 重点不是发明新 pattern，而是验证“FlashComm 整改是否可以走独立 pass 化路线”。",
        "复用 pattern 可以降低穿刺风险，也方便和现有 enable_sp 行为做对照。",
        "后续如果 SE 认可方向，可以再决定最终是保留 FlashCommDensePass，还是把命名和开关统一回 sequence_parallel。"
    ], 16)

    s = new_slide(prs, "验证方式")
    add_table(s, 0.55, 1.25, 12.2, 4.9, [
        ["验证层级", "当前/建议做法", "目的"],
        ["静态编译", "py_compile flashcomm_pass.py 和 graph_fusion_pass_manager.py", "保证代码语法与导入路径基本正确"],
        ["单测：范围判断", "is_applicable_for_range(start < / >= min_tokens)", "确认 token 阈值逻辑"],
        ["单测：pass 调用", "mock patterns.apply 返回 matched_count", "确认 pass 生命周期和 matched_count 记录"],
        ["单测：注册逻辑", "mock GraphFusionPassManager.configure", "确认 POC 开关和 enable_sp 互斥保护"],
        ["端到端 smoke", "dense TP>1 + enable_flashcomm_pass_poc", "确认真实图中 pattern 可命中并可运行"],
    ], 11)

    s = new_slide(prs, "当前方案的边界与风险")
    add_bullets(s, 0.75, 1.25, 12.0, 5.4, [
        "边界 1：POC 只证明 dense 基础通信/norm pattern 可被 pass 表达，不证明 MoE 叠加场景已解决。",
        "边界 2：由于和 enable_sp 互斥注册，这不是最终产品路径，只是避免穿刺阶段影响现有能力。",
        "边界 3：pad/unpad 仍在 custom op 和 runtime 路径中存在，POC 未解决“pad 上移到图外”的需求。",
        "风险：如果后续直接扩展到 MoE，多流和 shared_expert_dp 可能改变 FX 图中 pattern 的相邻关系和张量分布状态。",
    ], 16)

    s = new_slide(prs, "后续扩展路线建议")
    add_bullets(s, 0.75, 1.25, 12.0, 5.4, [
        "第一阶段：确认 dense pass POC 思路是否被 SE 接受，重点看 pattern 替换的语义是否正确。",
        "第二阶段：梳理 MoE 图中的 all_gather 后移 pattern，判断是新增 MoE pass，还是保留 prepare/finalize 特化逻辑。",
        "第三阶段：拆开 shared_expert_dp_enabled() 的混合语义，避免 FlashComm/SP 隐式改变 shared expert DP/TP 策略。",
        "第四阶段：单独设计多模态第一层和量化 allgather 位置，不把它们塞进基础 dense pass。",
    ], 16)

    s = new_slide(prs, "需要和 SE 讨论的问题")
    add_bullets(s, 0.75, 1.25, 12.0, 5.4, [
        "最终接口命名：继续叫 FlashComm，还是完全对齐主社区 sequence_parallel？",
        "基础 dense pass 是否可以作为后续整改主线，逐步替换 layer/custom op 中的显式分支？",
        "MoE 下 all_gather 后移是通过更多 FX pass 匹配，还是保留 runtime prepare/finalize 逻辑更稳？",
        "shared_expert_dp 和 multistream_moe 叠加时，哪些是正确性约束，哪些只是性能策略？",
        "pad/unpad 从 custom op 上移到 attention metadata 后，对 pass pattern 的输入输出 shape 约束是什么？",
    ], 16)

    s = new_slide(prs, "结论")
    add_bullets(s, 0.75, 1.45, 11.9, 4.7, [
        "这个 POC 是一个最小闭环：独立开关、独立 pass 名称、复用成熟 dense pattern、避免影响现有 enable_sp 路径。",
        "它适合用来向 SE 证明：FlashComm v1 的基础替换可以从 layer if/else 转向 FX pass。",
        "它不试图一次解决复杂叠加场景；复杂点应在基础方案确认后，通过 MoE、多流、shared expert、量化专题逐个扩展。",
    ], 18)

    prs.save(OUT)


if __name__ == "__main__":
    main()
