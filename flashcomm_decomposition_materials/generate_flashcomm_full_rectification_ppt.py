from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


TITLE_COLOR = RGBColor(31, 78, 121)
TEXT_COLOR = RGBColor(35, 35, 35)
ACCENT = RGBColor(91, 155, 213)
HEADER = RGBColor(221, 235, 247)
WARN = RGBColor(255, 242, 204)
OK = RGBColor(226, 239, 218)
FONT = "Microsoft YaHei"
OUT = "flashcomm_full_rectification_plan.pptx"


def setup():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    return prs


def set_font(paragraph, size, bold=False, color=TEXT_COLOR):
    paragraph.font.name = FONT
    paragraph.font.size = Pt(size)
    paragraph.font.bold = bold
    paragraph.font.color.rgb = color


def add_title(slide, title, subtitle=None):
    box = slide.shapes.add_textbox(Inches(0.45), Inches(0.25), Inches(12.4), Inches(0.6))
    p = box.text_frame.paragraphs[0]
    p.text = title
    set_font(p, 25, True, TITLE_COLOR)
    if subtitle:
        sub = slide.shapes.add_textbox(Inches(0.48), Inches(0.86), Inches(12.2), Inches(0.35))
        sp = sub.text_frame.paragraphs[0]
        sp.text = subtitle
        set_font(sp, 12, False, RGBColor(100, 100, 100))


def add_bullets(slide, left, top, width, height, bullets, font_size=15):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.08)
    tf.margin_right = Inches(0.08)
    for i, item in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        text, level = item if isinstance(item, tuple) else (item, 0)
        p.text = text
        p.level = level
        set_font(p, font_size if level == 0 else font_size - 1)
        p.space_after = Pt(4)


def add_table(slide, left, top, width, height, data, font_size=10):
    table = slide.shapes.add_table(len(data), len(data[0]), Inches(left), Inches(top), Inches(width), Inches(height)).table
    for r, row in enumerate(data):
        for c, text in enumerate(row):
            cell = table.cell(r, c)
            cell.text = text
            if r == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = HEADER
            for p in cell.text_frame.paragraphs:
                set_font(p, font_size, r == 0)
    return table


def add_box(slide, left, top, width, height, text, fill=RGBColor(242, 248, 252), font_size=13):
    shape = slide.shapes.add_shape(1, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = ACCENT
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = PP_ALIGN.CENTER
    set_font(p, font_size, True)
    return shape


def add_arrow(slide, left, top, width=0.55, height=0.3):
    shape = slide.shapes.add_shape(33, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT
    shape.line.color.rgb = ACCENT


def new_slide(prs, title, subtitle=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, title, subtitle)
    return slide


def main():
    prs = setup()

    s = new_slide(prs, "FlashComm v1 全特性耦合整改方案", "目标：以 FX pass 机制为主线，逐步日落 layer/custom op 中的显式分支")
    add_bullets(s, 0.75, 1.35, 11.9, 5.3, [
        "整改主线：把“通信形态替换”从 layer if/else 和自定义算子分支迁移到 FX graph pass。",
        "设计原则：先固化 dense 基线，再分层处理 MoE、shared_expert_dp、多流、多模态、量化和 padding。",
        "关键判断：不是所有耦合都适合 pass 化；运行时调度、metadata、第一层输入布局仍需要 runtime/metadata 层整改。",
        "最终目标：形成“pass 可表达的 pattern 库 + runtime 只负责布局/调度/metadata”的清晰边界。",
    ], 18)

    s = new_slide(prs, "当前问题总览")
    add_table(s, 0.45, 1.15, 12.45, 5.55, [
        ["问题域", "当前现状", "主要风险"],
        ["命名与开关", "FlashComm1、enable_sp、pass enable_sp、FlashCommDensePass POC 并存", "用户配置和运行时语义不一致"],
        ["padding", "maybe_all_gather / maybe_pad_and_reduce 内部做 pad/unpad/slice", "ACL graph capture 后 host int slice 固化"],
        ["shared_expert_dp", "SED 会强制 enable_sp；enable_sp 也影响 shared_expert_dp_enabled()", "显式配置与实际 shared expert 行为混淆"],
        ["multistream_moe", "多流是 stream/context 调度，但 shared_out 聚合受 SED/SP 影响", "单开多流语义不透明"],
        ["多模态", "第一层缺少 dense LLM 的 embedding 后通信前提", "第一层 pass pattern 不匹配"],
        ["量化", "dense / MoE / SP 中 quant 与 AG/RS 相对位置不统一", "通信带宽、精度和 pattern 冲突"],
    ], 10)

    s = new_slide(prs, "目标架构：分层收敛")
    add_box(s, 0.6, 1.25, 2.5, 0.8, "配置层\n统一接口与策略")
    add_arrow(s, 3.25, 1.48)
    add_box(s, 4.0, 1.25, 2.5, 0.8, "Metadata 层\npad/layout 上移")
    add_arrow(s, 6.65, 1.48)
    add_box(s, 7.4, 1.25, 2.5, 0.8, "FX Pass 层\npattern 替换")
    add_arrow(s, 10.05, 1.48)
    add_box(s, 10.8, 1.25, 2.2, 0.8, "Runtime 层\n调度/通信组")
    add_bullets(s, 0.75, 3.0, 12.0, 3.5, [
        "配置层：对齐社区 sequence_parallel 命名，FlashComm 作为 Ascend 内部实现名逐步弱化。",
        "Metadata 层：统一 token padding、actual length、first-layer layout，避免图内动态 slice。",
        "FX Pass 层：只做稳定算子 pattern 替换，如 AR+Norm、Norm+Quant+AG、MoE AG epilogue。",
        "Runtime 层：保留 stream/event、EP/DP/TP 组、MoE dispatch/combine、FC3 context 等图外调度职责。",
    ], 16)

    s = new_slide(prs, "整改边界：什么适合 pass，什么不适合")
    add_table(s, 0.55, 1.2, 12.2, 5.3, [
        ["类别", "适合程度", "建议"],
        ["all_reduce + RMSNorm", "适合 pass", "作为 dense 基线，已由 POC/SequenceParallelismPass 证明"],
        ["RMSNorm + quant + all_gather", "适合 pass", "扩展 norm_quant_fusion_pass，明确 AG 在 quant 后"],
        ["MoE all_gather 后移", "部分适合 pass", "用 SequenceParallelismMoePass 承接可稳定匹配的 epilogue"],
        ["padding / unpad / host num_tokens slice", "不适合图内处理", "上移到 attention metadata / runner 图外"],
        ["multistream gate/shared expert", "不适合纯 pass", "保留 runtime stream 调度，pass 只处理其中稳定算子片段"],
        ["多模态第一层输入布局", "不适合直接 pass", "先定义 metadata/layout，再决定是否新增 FirstLayerVL pattern"],
    ], 11)

    s = new_slide(prs, "阶段 0：统一术语和开关")
    add_table(s, 0.55, 1.25, 12.2, 4.7, [
        ["现有概念", "建议收敛方向", "备注"],
        ["VLLM_ASCEND_ENABLE_FLASHCOMM1", "迁移到 additional_config / pass_config 策略", "环境变量保留兼容期"],
        ["enable_sp()", "明确表示 runtime FC1/SP 状态", "不要同时承担 SED 判断"],
        ["pass_config.enable_sp", "对齐社区 sequence_parallel", "作为正式 pass 入口候选"],
        ["enable_flashcomm_pass_poc", "仅保留为穿刺开关", "验证完成后应删除或合并"],
        ["shared_expert_dp_enabled()", "拆成 explicit_sed / sp_enabled / pass_sp_enabled", "避免隐式改变 shared expert 语义"],
    ], 11)
    add_bullets(s, 0.75, 6.15, 11.8, 0.6, ["SE 决策点：最终外部接口是否完全对齐主社区 sequence_parallel。"], 14)

    s = new_slide(prs, "阶段 1：Dense 基线 pass 化")
    add_bullets(s, 0.75, 1.25, 12.0, 2.1, [
        "目标：先把最稳定的 dense FlashComm pattern 迁到 pass，证明替换语义和性能收益。",
        "基础 pattern：all_reduce + npu_add_rms_norm_bias -> reduce_scatter + npu_add_rms_norm_bias + all_gather。",
        "当前 POC：FlashCommDensePass 复用 Middle/LastAllReduceRMSNormPattern，并通过独立开关接入。",
    ], 16)
    add_table(s, 0.75, 3.85, 11.8, 2.1, [
        ["动作", "建议"],
        ["短期", "继续用 POC 向 SE 验证思路，补真实 dense TP>1 smoke"],
        ["中期", "合并到正式 SequenceParallelismPass 命名体系"],
        ["长期", "删除 linear_op 中可由 pass 表达的 FC1 特判"],
    ], 12)

    s = new_slide(prs, "阶段 2：padding 上移是前置整改")
    add_bullets(s, 0.75, 1.25, 12.0, 5.4, [
        "需求中最明确的工程问题：FlashComm 自定义算子里做 pad/unpad/slice 与图模式天然冲突。",
        "整改原则：模型图内不再根据 host 侧 num_tokens 做 slice；pad 规则统一在 attention metadata / runner 构造阶段完成。",
        "涉及路径：register_custom_ops.py 的 maybe_all_gather/maybe_pad_and_reduce/maybe_chunk_residual，distributed/utils.py 的 FC3 allgather/unpad。",
        "收益：pass pattern 的输入输出 shape 更稳定，多模态、MoE、量化都能共享统一 token layout 语义。",
    ], 16)

    s = new_slide(prs, "阶段 3：shared_expert_dp 解耦")
    add_table(s, 0.55, 1.2, 12.2, 5.25, [
        ["当前耦合", "问题", "整改建议"],
        ["SED 开启强制 enable_sp()", "SED 与 FlashComm/SP 绑定不透明", "保留正确性绑定，但以显式策略记录，不通过 helper 隐式修改"],
        ["shared_expert_dp_enabled() 包含 enable_sp()", "FlashComm 开启会影响 shared expert 是否 allreduce", "拆 helper，区分 explicit_sed 和 sp_runtime_state"],
        ["shared expert 命中 helper 后禁用 TP op", "显存劣化来源不直观", "文档和日志明确提示权重复制/显存代价"],
        ["MoE prepare/finalize 受 SED/replace_allreduce 共同影响", "张量 token 分布状态难追踪", "引入 layout state 注释或 metadata 标记"],
    ], 10)

    s = new_slide(prs, "阶段 4：multistream_moe 保留 runtime 调度边界")
    add_bullets(s, 0.75, 1.25, 12.0, 5.4, [
        "多流 MoE 的本质是调度优化：gate_stream、shared_expert_stream、quant_stream、FlashCommon3Context。",
        "这类 stream/event/context 依赖不适合全部放进 FX pass；pass 只处理其中稳定的算子子图。",
        "整改重点不是“给多流写一个大 pass”，而是明确 shared_out/topk/hidden_states 的布局状态。",
        "建议建立组合矩阵：SED on/off、SP on/off、multistream gate/shared、MoECommType、量化类型。",
    ], 16)

    s = new_slide(prs, "阶段 5：MoE 的 pass 化策略")
    add_table(s, 0.55, 1.2, 12.2, 5.25, [
        ["MoE 子问题", "建议机制", "理由"],
        ["标准 routed expert dispatch/combine", "runtime 保留", "依赖 EP/DP/通信组与 kernel"],
        ["MoE 前后 AG/RS epilogue", "pass + runtime 协同", "图上可匹配部分交给 SequenceParallelismMoePass"],
        ["shared expert 输出 allreduce", "runtime 保留并解耦 helper", "和 SED/TP 权重部署强相关"],
        ["AllGather+EP 下 W8A8 预量化", "runtime 先保留，后续验证可 pass 化", "与 quant_type/pertoken_scale 强相关"],
        ["MXFP8 prepare TODO", "专项决策", "目前与 W8A8 策略不一致"],
    ], 10)

    s = new_slide(prs, "阶段 6：多模态模型整改方案")
    add_bullets(s, 0.75, 1.25, 12.0, 5.4, [
        "核心问题：多模态第一层缺少 dense LLM 的 embedding 后通信前提，第一层 RMSNorm 前不是同构 pattern。",
        "短期：明确跳过 VL 第一层 dense FlashComm pattern，保留现有 layer0 runtime 特判，避免错误匹配。",
        "中期：把文本 embedding、vision merge、placeholder 写回后的 layout 用 metadata 显式描述。",
        "长期：如果数学上允许，把第一层构造成可匹配子图；否则新增 FirstLayerVLPattern，只处理明确模型族。",
        "deepstack 中间层已有 pass 思路：chunk(deepstack_input_embeds) 与 reduce_scatter 对齐，可作为中间层模板。",
    ], 16)

    s = new_slide(prs, "阶段 7：量化场景整改方案")
    add_table(s, 0.55, 1.15, 12.2, 5.55, [
        ["场景", "当前现状", "整改建议"],
        ["Dense norm+quant", "norm_quant_fusion_pass 已有 SP pattern：norm -> AG -> quant 替换为 norm_quant -> AG", "正式纳入 pass 化主线，验证 AG 在 quant 后的收益与正确性"],
        ["Row W8A8 + RS", "linear_op 中 quantize 后 npu_mm_reduce_scatter_base", "先保留 runtime，后续考虑 GEMM+RS pass/fusion"],
        ["MoE W8A8 prepare", "npu_dynamic_quant 在 EP allgather 前", "保留并补矩阵测试，确认 pertoken_scale AG 语义"],
        ["MoE MXFP8 prepare", "TODO：不预量化，留在 MoE MLP", "SE 决策：是否对齐 W8A8 或明确不支持"],
        ["multistream gate + quant", "W8A8 dynamic 从 FC3 context 取 topk", "保留 runtime context，pass 不表达 stream 依赖"],
    ], 10)

    s = new_slide(prs, "推荐 pass 库分层")
    add_table(s, 0.55, 1.2, 12.2, 5.2, [
        ["Pass 层", "负责 pattern", "备注"],
        ["DenseSequenceParallelPass", "AR + Norm -> RS + Norm + AG", "主社区 sequence_parallel 对齐"],
        ["NormQuantSPPass", "Norm + AG + Quant -> NormQuant + AG", "确保量化后通信"],
        ["MoeSequenceParallelPass", "MoE allgather epilogue / allgather+chunk no-op", "只处理稳定 epilogue"],
        ["VLMiddleLayerPass", "deepstack add + norm 的中间层 pattern", "多模态中间层，不处理第一层"],
        ["FirstLayerVLPass", "可选，模型族白名单", "只有当第一层 layout 被 metadata 固化后再做"],
    ], 11)

    s = new_slide(prs, "日落清单：从哪里开始删 if/else")
    add_table(s, 0.55, 1.15, 12.2, 5.55, [
        ["优先级", "日落对象", "前置条件"],
        ["P0", "模型图内 pad/unpad/slice", "attention metadata / runner 已统一 padding"],
        ["P1", "dense linear 中可由 AR+Norm pass 表达的 FC1 分支", "dense pass e2e 与性能通过"],
        ["P1", "env 与 pass 双轨同步胶水", "统一配置入口"],
        ["P2", "多模态中间层 deepstack runtime chunk 特判", "VL middle pass 稳定覆盖"],
        ["P2", "norm_quant 中重复 AG/quant 分支", "NormQuantSPPass 覆盖主量化类型"],
        ["P3", "MoE prepare/finalize 中可稳定图化的 epilogue", "MoE pass 与通信矩阵验证完成"],
    ], 10)

    s = new_slide(prs, "测试矩阵建议")
    add_table(s, 0.45, 1.05, 12.45, 5.8, [
        ["维度", "建议覆盖"],
        ["模型类型", "dense LLM、MoE、shared expert MoE、VL、VL+MoE 如有"],
        ["并行", "TP=1/2/4，EP on/off，DP on/off，prefill/decode"],
        ["开关", "sequence_parallel/pass、runtime FC1、SED、multistream_gate/shared"],
        ["通信类型", "ALLGATHER、ALLTOALL、MC2、FUSED_MC2"],
        ["量化", "BF16、W8A8、MXFP8、动态量化"],
        ["图模式", "eager、compile、ACL graph capture 不同 capture size"],
        ["正确性指标", "输出一致性、shape/layout invariants、matched_count、无图内动态 slice"],
        ["性能指标", "prefill 大 batch、decode 小 batch、通信量、显存峰值"],
    ], 10)

    s = new_slide(prs, "分阶段落地计划")
    add_table(s, 0.55, 1.15, 12.2, 5.55, [
        ["阶段", "目标", "交付物"],
        ["S0", "术语/开关对齐", "设计文档、配置迁移方案、helper 拆分方案"],
        ["S1", "Dense pass 正式化", "SequenceParallel pass 路线确认，POC 合并/日落"],
        ["S2", "Padding 上移", "删除 custom op 内动态 slice，metadata 统一 pad"],
        ["S3", "量化 pass 扩展", "NormQuantSPPass 覆盖 W8A8 dense，明确 MXFP8 策略"],
        ["S4", "MoE/SED 解耦", "shared_expert_dp_enabled 拆分，MoE 通信矩阵测试"],
        ["S5", "多模态/VL", "第一层 layout 方案，中间层 pass 覆盖，模型白名单"],
        ["S6", "多流收敛", "runtime 调度边界文档，FC3 allgather/pad 清理"],
    ], 10)

    s = new_slide(prs, "需要 SE 决策的问题")
    add_bullets(s, 0.75, 1.25, 12.0, 5.4, [
        "外部接口是否完全对齐主社区 sequence_parallel，FlashComm 是否只作为 Ascend 内部实现名。",
        "SED 与 FlashComm/SP 的绑定是否是正确性要求；如果只是性能策略，是否允许用户显式选择。",
        "多流 MoE 在 SED 关闭、SP 关闭时是否允许单开；允许的话支持哪些 MoECommType 和模型白名单。",
        "多模态第一层是选择构造同构通信子图，还是长期用 FirstLayerVL 特化规则。",
        "MXFP8 是否要对齐 W8A8 的“量化后 AG”策略，还是明确不支持该优化。",
        "padding 上移方案由哪个需求统一承接，完成前哪些 pass 只能作为实验能力。",
    ], 16)

    s = new_slide(prs, "建议结论")
    add_bullets(s, 0.75, 1.35, 11.9, 5.3, [
        "完整整改不应追求“一把大 pass 覆盖所有场景”，而应分成 stable pattern pass、metadata 统一、runtime 调度三条线。",
        "Dense + 量化的稳定 pattern 优先 pass 化；MoE/多流/shared expert 先解耦语义，再逐步把 epilogue 图化。",
        "多模态第一层是结构性例外，必须先解决输入 layout，再决定是否新增专用 pass。",
        "padding 上移是全局前置条件，它决定后续 pass pattern 是否能稳定且图模式安全。",
        "建议用现有 FlashCommDensePass POC 作为 SE 讨论的最小样例，但最终收敛到社区 sequence_parallel 命名和 pass 体系。",
    ], 18)

    prs.save(OUT)


if __name__ == "__main__":
    main()
