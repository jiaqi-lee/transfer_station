from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


TITLE_COLOR = RGBColor(31, 78, 121)
TEXT_COLOR = RGBColor(35, 35, 35)
ACCENT = RGBColor(91, 155, 213)
HEADER = RGBColor(221, 235, 247)
FONT = "Microsoft YaHei"


def setup():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    return prs


def set_font(paragraph, size, bold=False, color=TEXT_COLOR):
    paragraph.font.name = FONT
    paragraph.font.size = Pt(size)
    paragraph.font.bold = bold
    paragraph.font.color.rgb = color


def add_title(slide, title, subtitle=None):
    box = slide.shapes.add_textbox(Inches(0.45), Inches(0.25), Inches(12.4), Inches(0.6))
    p = box.text_frame.paragraphs[0]
    p.text = title
    set_font(p, 26, True, TITLE_COLOR)
    if subtitle:
        sub = slide.shapes.add_textbox(Inches(0.48), Inches(0.86), Inches(12.2), Inches(0.35))
        sp = sub.text_frame.paragraphs[0]
        sp.text = subtitle
        set_font(sp, 12, False, RGBColor(100, 100, 100))


def add_bullets(slide, left, top, width, height, bullets, font_size=15):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.08)
    tf.margin_right = Inches(0.08)
    for i, item in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        text, level = item if isinstance(item, tuple) else (item, 0)
        p.text = text
        p.level = level
        set_font(p, font_size if level == 0 else font_size - 1)
        p.space_after = Pt(5)


def add_label(slide, left, top, text):
    shape = slide.shapes.add_shape(1, Inches(left), Inches(top), Inches(2.8), Inches(0.36))
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT
    shape.line.color.rgb = ACCENT
    p = shape.text_frame.paragraphs[0]
    p.text = text
    p.alignment = PP_ALIGN.CENTER
    set_font(p, 12, True, RGBColor(255, 255, 255))


def add_table(slide, left, top, width, height, data, font_size=11):
    table = slide.shapes.add_table(len(data), len(data[0]), Inches(left), Inches(top), Inches(width), Inches(height)).table
    for r, row in enumerate(data):
        for c, text in enumerate(row):
            cell = table.cell(r, c)
            cell.text = text
            if r == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = HEADER
            for p in cell.text_frame.paragraphs:
                set_font(p, font_size, r == 0)


def new_slide(prs, title, subtitle=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, title, subtitle)
    return slide


def make_sed():
    prs = setup()

    s = new_slide(prs, "shared_expert_dp 与 FlashComm 耦合梳理", "面向 FlashComm v1 pass 化整改的背景材料")
    add_bullets(s, 0.75, 1.45, 11.9, 4.7, [
        "目标：讲清 shared_expert_dp 为什么会牵动 FlashComm/SP，以及当前耦合点在哪里。",
        "结论先行：enable_shared_expert_dp 开启时会强制内部 enable_sp() 为 True；反过来，FlashComm/SP 开启时也会让 shared_expert_dp_enabled() 返回 True。",
        "风险：配置语义、运行时张量分布状态、shared expert 是否 TP all-reduce 被一个混合 helper 绑定，导致特性叠加不透明。",
        "建议：和 SE 对齐时把“用户显式 SED”“sequence parallel/flashcomm 状态”“pass-based SP 状态”拆开讨论。",
    ], 18)

    s = new_slide(prs, "术语与问题边界")
    add_label(s, 0.7, 1.25, "本页结论")
    add_bullets(s, 0.75, 1.75, 5.8, 4.9, [
        "shared expert：MoE 中所有 token 都会经过的共享专家，与 routed experts 并列。",
        "TP 切分 shared expert：权重/输出需要 TP 通信合并。",
        "shared_expert_dp：当前 Ascend 中更偏向复制 shared expert，减少 shared expert 侧通信/重复工作。",
        "FlashComm1 / enable_sp：把部分 TP all-reduce 改成 reduce-scatter + 局部 norm + all-gather。",
    ])
    add_bullets(s, 6.9, 1.75, 5.7, 4.9, [
        "要回答的问题：",
        ("为什么 SED 会默认拉起 FlashComm/SP？", 1),
        ("为什么 FlashComm 开启时不会默认打开配置项 enable_shared_expert_dp？", 1),
        ("哪些分支实际受 shared_expert_dp_enabled() 影响？", 1),
        ("这种耦合是正确性、性能还是历史实现耦合？", 1),
    ])

    s = new_slide(prs, "开关链路：配置项和 helper 语义不一致")
    add_table(s, 0.55, 1.35, 12.2, 4.1, [
        ["对象", "位置", "语义"],
        ["enable_shared_expert_dp", "ascend_config.py", "additional_config 为 true 且 EP 开启且 TP>1 时才真正为 true"],
        ["enable_sp()", "utils.py", "读取 FlashComm1 env；若 enable_shared_expert_dp=True 且原本未开，会把 _ENABLE_SP 置 True"],
        ["shared_expert_dp_enabled()", "utils.py", "enable_shared_expert_dp or enable_sp() or enable_sp_by_pass()"],
        ["flash_comm_v1_enabled", "ascend_forward_context.py", "MoE 下 enable_sp 且 num_tokens 非 None；dense 下额外要求 num_tokens>1000"],
    ], 12)
    add_bullets(s, 0.75, 5.75, 11.9, 1.0, [
        "关键点：enable_shared_expert_dp 是用户显式配置语义；shared_expert_dp_enabled() 是运行时混合语义，包含 FlashComm/SP。"
    ], 15)

    s = new_slide(prs, "SED 开启为什么会拉起 FlashComm/SP")
    add_bullets(s, 0.75, 1.35, 12.0, 4.9, [
        "ascend_config.py 中 enable_shared_expert_dp 真正生效后，会 assert enable_sp(vllm_config, enable_shared_expert_dp=True)。",
        "utils.enable_sp() 如果没有通过环境变量打开 FlashComm1，但发现 enable_shared_expert_dp=True，会把内部 _ENABLE_SP 置 True。",
        "这会让 MoE forward context 的 flash_comm_v1_enabled 进入 true 分支，并影响 MoE prepare/finalize 的 replace_allreduce。",
        "所以“SED 开启默认开启 FlashComm”更精确的表述是：SED 开启会强制内部 SP 标志为 True，不是改环境变量。",
    ], 16)

    s = new_slide(prs, "FlashComm 开启为什么不会默认开启 SED 配置项")
    add_bullets(s, 0.75, 1.25, 12.0, 5.3, [
        "enable_shared_expert_dp 会导致 shared expert 更偏复制，显存更差；文档明确“性能更好但占用更多显存”。",
        "FlashComm/SP 是通信与激活切分优化，默认不应改变用户的 shared expert 部署策略。",
        "但当前 helper 把 enable_sp() 算入 shared_expert_dp_enabled()，导致 FlashComm 开启后部分 shared expert 行为仍会按 SED-like 语义执行。",
        "这就是当前方案最需要和 SE 对齐的点：配置项没有自动打开，但运行时行为已经被 SP 间接改变。",
    ], 16)

    s = new_slide(prs, "运行时影响点 1：shared expert 是否 TP all-reduce")
    add_bullets(s, 0.75, 1.25, 12.0, 5.2, [
        "fused_moe.py 中，ALLTOALL / MC2 / FUSED_MC2 通信类型下，只有 not shared_expert_dp_enabled() 时才对 shared_out 做 tensor_model_parallel_all_reduce。",
        "因此只要 FlashComm1/SP 或 enable_sp_by_pass 打开，即使用户没有显式设置 enable_shared_expert_dp，shared_out 的 TP all-reduce 分支也会被跳过。",
        "这会改变 shared expert 输出的聚合时机，需要和 routed expert finalize 的通信状态一起理解。",
    ], 16)

    s = new_slide(prs, "运行时影响点 2：MoE prepare/finalize 的切分与恢复")
    add_table(s, 0.55, 1.25, 12.2, 4.8, [
        ["路径", "未开 SED/FC", "开 SED 或 FC 相关路径"],
        ["All2All prepare", "TP 维 pad 后按 token 切分", "replace_allreduce 或 SED 下跳过 TP pad/slice"],
        ["All2All finalize", "TP all-gather 恢复并 unpad", "replace_allreduce 或 SED 下跳过 gather/unpad"],
        ["MC2 prepare", "使用 mc2_mask，必要时 pad/slice", "SED 下跳过 pad/slice；replace_allreduce 下跳过主体逻辑"],
        ["AllGather+EP", "DP AG -> MoE -> DP RS", "enable_sp/enable_sp_by_pass 下走 EP AG/RS 优化路径"],
    ], 11)
    add_bullets(s, 0.75, 6.25, 11.8, 0.55, ["核心：SED 与 FlashComm 都在改变 MoE 前后 hidden_states 的 token 分布状态。"], 14)

    s = new_slide(prs, "显存劣化来源")
    add_bullets(s, 0.75, 1.25, 12.0, 5.4, [
        "shared expert 权重：linear_op.py 中 shared_expert* 命中 shared_expert_dp_enabled() 时不挂 TP 自定义 op，等价更偏每卡完整 shared expert。",
        "激活与缓冲：跳过 TP token 切分后，单卡上的 token 维中间张量可能更大。",
        "通信缓冲：AllGather/ReduceScatter/FC3 上下文、多流事件与缓冲可能增加峰值，具体需要按模型和通信类型实测。",
        "所以 FlashComm 不默认打开 enable_shared_expert_dp 是合理的：它不是纯通信优化，而会影响权重/激活内存形态。",
    ], 16)

    s = new_slide(prs, "与上游 vLLM 的对照")
    add_table(s, 0.55, 1.25, 12.2, 4.9, [
        ["维度", "上游 vLLM", "vLLM Ascend 当前"],
        ["shared expert reduce", "runner 根据 kernel output 是否已 reduce 决定是否提前 reduce shared_out", "按 MoECommType + shared_expert_dp_enabled() 决定是否 TP all-reduce"],
        ["sequence parallel", "社区命名为 sequence_parallel，主要 dense pass", "FlashComm1/enable_sp 与 pass、custom op、MoE 路径混合存在"],
        ["SED 开关", "无同名 enable_shared_expert_dp", "Ascend additional_config，DeepSeek 等模型推荐手动开"],
        ["耦合风险", "reduce 策略相对集中", "SED、SP、pass-based SP 被 helper 混合"],
    ], 11)

    s = new_slide(prs, "当前组合矩阵")
    add_table(s, 0.55, 1.25, 12.2, 4.9, [
        ["显式 SED", "FlashComm/SP", "shared_expert_dp_enabled()", "当前行为解读"],
        ["关", "关", "False", "传统 shared expert TP 通信路径"],
        ["开", "原本关", "True", "SED 会强制 enable_sp=True，进入 SP/FC 相关路径"],
        ["关", "开", "True", "配置项未开，但 shared expert 分支按 SED-like 行为处理"],
        ["开", "开", "True", "推荐组合之一，但显存与通信形态都变化"],
    ], 12)

    s = new_slide(prs, "问题归类：正确性、性能、历史耦合")
    add_bullets(s, 0.75, 1.25, 12.0, 5.4, [
        "正确性耦合：MoE routed/shared 输出在 TP/DP/EP 维度的聚合状态必须一致，否则结果错误。",
        "性能耦合：SED 复制 shared expert 换通信减少，但显存增加；FlashComm/SP 通过 token shard 降低 norm/通信成本。",
        "历史耦合：shared_expert_dp_enabled() 同时包含显式 SED、enable_sp、enable_sp_by_pass，导致开关语义不清。",
        "整改方向：先保留正确性耦合，逐步把性能策略和历史胶水拆为显式配置或 pass 策略。",
    ], 16)

    s = new_slide(prs, "给 SE 的建议问题")
    add_bullets(s, 0.75, 1.25, 12.0, 5.4, [
        "是否认可把 enable_shared_expert_dp、enable_sp/flashcomm、enable_sp_by_pass 拆成三个独立语义？",
        "FlashComm/SP 开启时是否应该继续让 shared_expert_dp_enabled() 返回 True？如果是，命名是否需要调整？",
        "SED 与 FlashComm 的绑定是正确性要求，还是只是当前实现为了避免重复通信的策略？",
        "pass 化后，MoE prepare/finalize 里哪些分支应该保留，哪些可以被 pattern pass 替代？",
    ], 16)

    prs.save("shared_expert_dp_flashcomm_coupling.pptx")


def make_multistream():
    prs = setup()

    s = new_slide(prs, "multistream_moe 与 FlashComm / SED 叠加梳理", "面向 FlashComm v1 pass 化整改的背景材料")
    add_bullets(s, 0.75, 1.45, 11.9, 4.7, [
        "当前代码中 enable_multistream_moe 不是实际生效的主开关；生产逻辑主要看 multistream_overlap_gate 和 multistream_overlap_shared_expert。",
        "多流本质是调度/overlap 优化：把 shared expert 或 gate/topk 放到独立 NPU stream，与 routed experts 通信/计算重叠。",
        "它与 FlashComm/SED 的耦合来自 shared expert 输出聚合、EP all-gather、FC3 context 和 hidden state 的 token 分布状态。",
    ], 18)

    s = new_slide(prs, "开关现状：文档名与代码名不完全一致")
    add_table(s, 0.55, 1.25, 12.2, 4.7, [
        ["开关", "是否生产生效", "说明"],
        ["enable_multistream_moe", "未发现生产读取", "出现在部分测试/YAML，当前代码未作为实际控制项"],
        ["multistream_overlap_shared_expert", "生效", "shared expert 独立 stream，与 dispatch/combine 重叠"],
        ["multistream_overlap_gate", "生效", "gate stream + FlashCommon3Context，提前计算 shared_out/topk"],
        ["enable_shared_expert_dp", "间接影响", "影响 shared expert TP/DP 语义和 all-reduce 分支"],
    ], 12)
    add_bullets(s, 0.75, 6.1, 11.8, 0.6, ["建议汇报时把“多流 MoE”明确拆成两个实际开关，避免和 enable_multistream_moe 混淆。"], 14)

    s = new_slide(prs, "多流 shared expert 路径")
    add_bullets(s, 0.75, 1.25, 12.0, 5.4, [
        "AscendSharedFusedMoE._forward_shared_experts 中使用 shared_experts_calculation_stream()。",
        "shared expert 的 gate_up/activation/down_proj 分段执行，分别等待 routed experts 的 before_dispatch / before_combine event。",
        "目标是让 shared expert 计算与 routed experts 的 dispatch/combine 通信重叠。",
        "最后 default stream wait shared expert stream，保证 shared_out 可与 routed_out 合并。",
    ], 16)

    s = new_slide(prs, "多流 gate / FlashComm3 路径")
    add_bullets(s, 0.75, 1.25, 12.0, 5.4, [
        "multistream_overlap_gate 开启时，AscendSharedFusedMoE 先 set_flash_common3_context(shared_experts=...)。",
        "AscendFusedMoE.forward_impl 在 gate_stream 上执行 shared_experts(hidden_states) 和 select_experts。",
        "shared_out、topk_weights、topk_ids 通过 FlashCommon3Context 传回主流。",
        "PrepareAndFinalizeWithAllGather 在该场景使用 fc3_all_gather_and_maybe_unpad_impl，和普通 maybe_all_gather 路径不同。",
    ], 16)

    s = new_slide(prs, "与 shared_expert_dp 的耦合点")
    add_bullets(s, 0.75, 1.25, 12.0, 5.4, [
        "多流 shared expert/gate 都会产生 shared_out；shared_out 是否需要 TP all-reduce 仍由 shared_expert_dp_enabled() 判断。",
        "shared_expert_dp_enabled() 又包含 enable_sp() 和 enable_sp_by_pass()，所以 FlashComm/SP 会改变多流路径的 shared_out 聚合行为。",
        "因此多流本身没有硬依赖 SED，但 TP>1 + EP + shared expert 下，不开 SED 时语义更容易不透明。",
    ], 16)

    s = new_slide(prs, "与 FlashComm 的耦合点")
    add_table(s, 0.55, 1.25, 12.2, 4.9, [
        ["耦合点", "说明"],
        ["enable_sp / flash_comm_v1_enabled", "决定 MoE prepare 的 replace_allreduce 和若干 custom op 行为"],
        ["maybe_all_gather_and_maybe_unpad", "普通 all-gather + unpad 路径，包含 pad/slice 技术债"],
        ["fc3_all_gather_and_maybe_unpad_impl", "multistream_overlap_gate 下使用的 FC3 all-gather 路径"],
        ["maybe_pad_and_reduce", "SP finalize 后 reduce_scatter/pad 相关路径"],
        ["shared_expert_dp_enabled", "FlashComm/SP 间接影响 shared_out 是否 all-reduce"],
    ], 11)

    s = new_slide(prs, "当前组合矩阵")
    add_table(s, 0.45, 1.10, 12.45, 5.55, [
        ["SED", "FlashComm/SP", "multistream", "当前判断"],
        ["开", "开", "开", "可开，当前文档/模型配置常见；但显存和通信形态都变化"],
        ["开", "关", "开", "实际会被 SED 拉起 enable_sp，不能视为纯关 FlashComm"],
        ["关", "开", "开", "配置上可开，但 shared_expert_dp_enabled=True，行为不是纯 SED off"],
        ["关", "关", "开", "代码无硬禁止；TP>1+EP+shared expert 下需按模型实测正确性和 shape"],
        ["任意", "任意", "关", "回到普通 MoE shared/routed 顺序执行路径"],
    ], 11)

    s = new_slide(prs, "多流对 FX 图/pass 的潜在影响")
    add_bullets(s, 0.75, 1.25, 12.0, 5.4, [
        "多流引入 stream/event/context，不只是纯算子序列替换，FX pass 不一定天然表达调度语义。",
        "multistream_overlap_gate 把 topk/shared_out 提前到 gate_stream，可能改变图中 all_gather、topk、shared expert 的相邻关系。",
        "pass 化时应先处理纯 dense pattern；MoE 多流是否用更多 pass，取决于 FX 图中是否稳定出现可匹配 pattern。",
        "如果多流逻辑主要由 Python stream/context 控制，可能不适合完全用 FX pass 承担。",
    ], 16)

    s = new_slide(prs, "与量化的耦合")
    add_bullets(s, 0.75, 1.25, 12.0, 5.4, [
        "w8a8_dynamic.py 在 multistream_overlap_gate 下读取 FlashCommon3Context 中的 topk 信息。",
        "prepare_finalize.py AllGather+EP 路径对 W8A8 会先做 npu_dynamic_quant，再 all-gather pertoken_scale。",
        "需求里提到 dense 量化希望 quant 插在 allgather 前，这是 pass pattern 设计问题；MoE gate/topk 量化路径需要单独确认。",
    ], 16)

    s = new_slide(prs, "风险与建议拆解")
    add_bullets(s, 0.75, 1.25, 12.0, 5.4, [
        "风险 1：enable_multistream_moe 命名与实际开关不一致，排障容易误判。",
        "风险 2：多流是否能单开，不是配置层问题，而是 shared_out 聚合与 hidden state 分布状态是否一致的问题。",
        "风险 3：FC3 all-gather/unpad 仍有 pad/slice 技术债，和需求中“pad 上移到图外”冲突。",
        "建议：先把多流定义为调度优化，不让它隐式决定 shared expert TP/DP 语义。",
    ], 16)

    s = new_slide(prs, "给 SE 的建议问题")
    add_bullets(s, 0.75, 1.25, 12.0, 5.4, [
        "后续产品/代码中是否保留 enable_multistream_moe 这个别名，还是统一到 multistream_overlap_gate/shared_expert？",
        "TP>1 + EP + shared expert 下，SED 关闭时是否允许多流单开？需要哪些模型/通信类型白名单？",
        "FC3 all-gather/unpad 是否应该和 FlashComm1 的 pad 上移整改一起处理？",
        "多流叠加是否应该通过更多 FX pass 解决，还是保留 runtime stream 调度，只把通信/norm pattern pass 化？",
    ], 16)

    prs.save("multistream_moe_flashcomm_coupling.pptx")


if __name__ == "__main__":
    make_sed()
    make_multistream()
